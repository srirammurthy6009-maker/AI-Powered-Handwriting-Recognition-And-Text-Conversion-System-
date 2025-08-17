from flask import Flask, request, send_file, render_template_string
import easyocr
import numpy as np
import cv2
from PIL import Image
from docx import Document
from pptx import Presentation
import tempfile
import os
import openai
from pdf2image import convert_from_bytes

app = Flask(__name__)
UPLOAD_FOLDER = tempfile.gettempdir()
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# 🔐 Inbuilt OpenAI API Key
openai.api_key = "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"

HTML = '''
<!DOCTYPE html>
<html>
<head>
    <title>Text2Word AI+ - Flask</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 30px; background: #f2f2f2; }
        h2 { color: #333; }
        textarea { width: 100%; height: 200px; font-family: monospace; }
        input[type="file"], input[type="submit"] {
            padding: 10px; margin: 10px 0; display: block; width: 100%;
        }
        a { text-decoration: none; background: #28a745; color: white; padding: 10px 20px; display: inline-block; margin-top: 10px; }
        .error { color: red; }
    </style>
</head>
<body>
    <h2>📝 Text2Word AI+ - Handwriting to Word</h2>
    <form method="POST" enctype="multipart/form-data">
        <label>📤 Upload Image, PDF, or PPT:</label>
        <input type="file" name="file" required>
        <input type="submit" value="Extract and Clean Text">
    </form>

    {% if error %}
        <p class="error">❌ {{ error }}</p>
    {% endif %}

    {% if extracted_text %}
        <h3>📄 Raw Extracted Text:</h3>
        <textarea readonly>{{ extracted_text }}</textarea>
    {% endif %}

    {% if corrected_text %}
        <h3>🤖 GPT-Corrected Text:</h3>
        <textarea readonly>{{ corrected_text }}</textarea>
    {% endif %}

    {% if download_ready %}
        <a href="/download">📥 Download Word Document</a>
    {% endif %}
</body>
</html>
'''

def extract_text_from_image(file_stream):
    file_bytes = np.asarray(bytearray(file_stream.read()), dtype=np.uint8)
    image_np = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)

    # Improved handwriting preprocessing
    gray = cv2.cvtColor(image_np, cv2.COLOR_BGR2GRAY)
    blur = cv2.medianBlur(gray, 3)
    adapt_thresh = cv2.adaptiveThreshold(blur, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                         cv2.THRESH_BINARY_INV, 15, 10)

    reader = easyocr.Reader(['en'], gpu=False)
    result = reader.readtext(adapt_thresh, detail=0)
    return "\n".join(result)

def extract_text_from_pdf(file_stream):
    images = convert_from_bytes(file_stream.read())
    reader = easyocr.Reader(['en'], gpu=False)
    full_text = ""
    for img in images:
        img_np = np.array(img)
        result = reader.readtext(img_np, detail=0)
        full_text += "\n".join(result) + "\n\n"
    return full_text

def extract_text_from_ppt(file_stream):
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def correct_text_with_gpt(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Fix grammar, spelling, and punctuation:\n{text}"}],
            temperature=0.2
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"⚠ GPT Correction Failed: {str(e)}"

def save_to_word(text):
    doc = Document()
    doc.add_paragraph(text)
    path = os.path.join(tempfile.gettempdir(), "output_text.docx")
    doc.save(path)
    return path

@app.route("/", methods=["GET", "POST"])
def index():
    extracted_text = ""
    corrected_text = ""
    error = ""
    download_ready = False

    if request.method == "POST":
        file = request.files.get("file")

        if not file:
            error = "Please upload a file."
        else:
            filetype = file.content_type
            file.stream.seek(0)

            try:
                if "image" in filetype:
                    extracted_text = extract_text_from_image(file.stream)
                elif "pdf" in filetype:
                    file.stream.seek(0)
                    extracted_text = extract_text_from_pdf(file.stream)
                elif "presentation" in filetype:
                    extracted_text = extract_text_from_ppt(file.stream)
                else:
                    error = "Unsupported file type."

                corrected_text = correct_text_with_gpt(extracted_text)
                save_to_word(corrected_text)
                download_ready = True

            except Exception as e:
                error = str(e)

    return render_template_string(HTML,
                                  extracted_text=extracted_text,
                                  corrected_text=corrected_text,
                                  download_ready=download_ready,
                                  error=error)

@app.route("/download")
def download():
    path = os.path.join(tempfile.gettempdir(), "output_text.docx")
    return send_file(path, as_attachment=True, download_name="extracted_text.docx")

if __name__ == "__main__":
    app.run(debug=True, port=5000)
